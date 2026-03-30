import os
import pandas as pd
from datetime import datetime

# ─── MS Word ─────────────────────────────────────────────────────────────
try:
    from docx import Document
    from docx.shared import Pt, RGBColor, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH
except ImportError:
    pass

def create_word(filepath):
    doc = Document()
    
    # Title
    title = doc.add_heading('NEXUS Executive Summary Report', 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    for run in title.runs:
        run.font.color.rgb = RGBColor(41, 128, 185) # Blue
    
    doc.add_paragraph(f"Generated on: {datetime.now().strftime('%B %d, %Y')}")
    
    # Heading 1
    h1 = doc.add_heading('1. Project Overview', level=1)
    for run in h1.runs:
        run.font.color.rgb = RGBColor(52, 73, 94)
    
    p = doc.add_paragraph(
        "NEXUS is a next-generation local AI platform that provides autonomous "
        "desktop agents capable of creating, editing, and managing complex "
        "office documents seamlessly. This report demonstrates the power of "
        "automated document generation with rich styling and dynamic content."
    )
    
    # Table
    doc.add_heading('2. Performance Metrics', level=2)
    table = doc.add_table(rows=1, cols=3)
    table.style = 'Light Shading Accent 1'
    hdr_cells = table.rows[0].cells
    hdr_cells[0].text = 'Metric'
    hdr_cells[1].text = 'Q1 Target'
    hdr_cells[2].text = 'Q1 Actual'
    
    data = [
        ('Speed (ms)', '500', '420'),
        ('Accuracy (%)', '95.0', '98.5'),
        ('Tasks Completed', '10,000', '12,500')
    ]
    for m, t, a in data:
        row = table.add_row().cells
        row[0].text = m
        row[1].text = t
        row[2].text = a

    doc.add_page_break()
    doc.add_heading('3. Conclusion', level=1)
    doc.add_paragraph("Automation significantly improves desktop productivity.")
    
    doc.save(filepath)
    print(f"Created: {filepath}")

# ─── Excel ───────────────────────────────────────────────────────────────
try:
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils.dataframe import dataframe_to_rows
    from openpyxl.chart import BarChart, Reference
except ImportError:
    pass

def create_excel(filepath):
    wb = Workbook()
    ws = wb.active
    ws.title = "Sales Analytics"
    
    # Header Styling
    header_fill = PatternFill(start_color="2980B9", end_color="2980B9", fill_type="solid")
    header_font = Font(color="FFFFFF", bold=True, size=12)
    center_align = Alignment(horizontal="center", vertical="center")
    thin_border = Border(left=Side(style='thin'), right=Side(style='thin'),
                         top=Side(style='thin'), bottom=Side(style='thin'))
    
    headers = ['Month', 'Revenue ($)', 'Profit ($)', 'Growth (%)']
    ws.append(headers)
    for col in range(1, 5):
        cell = ws.cell(row=1, column=col)
        cell.fill = header_fill
        cell.font = header_font
        cell.alignment = center_align
        cell.border = thin_border
        
    data = [
        ['Jan', 45000, 15000, 5.2],
        ['Feb', 52000, 18000, 15.5],
        ['Mar', 48000, 14000, -7.6],
        ['Apr', 61000, 22000, 27.0],
        ['May', 75000, 28000, 22.9]
    ]
    
    for row_idx, row_data in enumerate(data, 2):
        for col_idx, val in enumerate(row_data, 1):
            cell = ws.cell(row=row_idx, column=col_idx, value=val)
            cell.alignment = center_align
            cell.border = thin_border
            if col_idx == 4:
                 cell.number_format = '0.0%'
    
    # Adjust column widths
    ws.column_dimensions['A'].width = 15
    ws.column_dimensions['B'].width = 20
    ws.column_dimensions['C'].width = 20
    ws.column_dimensions['D'].width = 15
    
    # Create Chart
    chart = BarChart()
    chart.title = "Revenue vs Profit (Q1-Q2)"
    chart.style = 13
    chart.x_axis.title = 'Month'
    chart.y_axis.title = 'Amount ($)'
    
    data_ref = Reference(ws, min_col=2, min_row=1, max_row=6, max_col=3)
    cats = Reference(ws, min_col=1, min_row=2, max_row=6)
    chart.add_data(data_ref, titles_from_data=True)
    chart.set_categories(cats)
    
    ws.add_chart(chart, "F2")
    
    wb.save(filepath)
    print(f"Created: {filepath}")

# ─── PowerPoint ──────────────────────────────────────────────────────────
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor as PPTX_RGBColor
except ImportError:
    pass

def create_ppt(filepath):
    prs = Presentation()
    
    # Slide 1: Title
    title_slide_layout = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(title_slide_layout)
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    
    title.text = "NEXUS Agent Overview"
    title.text_frame.paragraphs[0].font.color.rgb = PPTX_RGBColor(41, 128, 185)
    title.text_frame.paragraphs[0].font.bold = True
    subtitle.text = "Autonomous AI operations on the Desktop"
    
    # Slide 2: Content
    bullet_slide_layout = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(bullet_slide_layout)
    title2 = slide2.shapes.title
    body2 = slide2.placeholders[1]
    
    title2.text = "Key Capabilities"
    tf = body2.text_frame
    
    p = tf.add_paragraph()
    p.text = "Local Model Execution"
    p.level = 0
    p.font.size = Pt(24)
    
    p = tf.add_paragraph()
    p.text = "No internet required for processing"
    p.level = 1
    p.font.color.rgb = PPTX_RGBColor(100, 100, 100)
    
    p = tf.add_paragraph()
    p.text = "Office Automation"
    p.level = 0
    p.font.size = Pt(24)
    
    p = tf.add_paragraph()
    p.text = "Excel, Word, PowerPoint manipulation"
    p.level = 1
    p.font.color.rgb = PPTX_RGBColor(100, 100, 100)
    
    prs.save(filepath)
    print(f"Created: {filepath}")


if __name__ == "__main__":
    out_dir = os.path.dirname(os.path.abspath(__file__))
    word_path = os.path.join(out_dir, "NEXUS_Report_Stylized.docx")
    excel_path = os.path.join(out_dir, "Sales_Analytics_Stylized.xlsx")
    ppt_path = os.path.join(out_dir, "NEXUS_Presentation_Stylized.pptx")
    
    create_word(word_path)
    create_excel(excel_path)
    create_ppt(ppt_path)
    print("All impressive documents generated successfully!")
